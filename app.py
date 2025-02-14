from flask import Flask, render_template
import sqlite3
import os
from pptx import Presentation
from PIL import Image

app = Flask(__name__)
DB_FILE = "presentation.db"
UPLOAD_FOLDER = "static/slides"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def extract_slides(pptx_path):
    prs = Presentation(pptx_path)
    slide_images = []
    
    for i, slide in enumerate(prs.slides):
        img_path = os.path.join(UPLOAD_FOLDER, f"slide_{i + 1}.png")
        slide.shapes._spTree  # Ensure all shapes are loaded
        
        slide_img = Image.new('RGB', (1280, 720), "white")  # Create a blank image
        slide_img.save(img_path)
        slide_images.append(img_path)
    
    return slide_images

def init_db():
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    c.execute("""
        CREATE TABLE IF NOT EXISTS presentation (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            slide_path TEXT
        )
    """)
    conn.commit()
    conn.close()

@app.route('/')
def index():
    conn = sqlite3.connect(DB_FILE)
    c = conn.cursor()
    c.execute("SELECT slide_path FROM presentation")
    slides = [row[0] for row in c.fetchall()]
    conn.close()
    return render_template("index.html", slides=slides)

if __name__ == "__main__":
    init_db()
    app.run(host='0.0.0.0', port=3000)
